from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

def create_presentation():
    # 1. 젠테이션 객체 생성
    prs = Presentation()

    # 2. 슬라이드 레이아웃 설정 (0: 제목 슬라이드, 1: 제목 및 내용)
    title_slide_layout = prs.slide_layouts[0]
    bullet_slide_layout = prs.slide_layouts[1]

    # --- Slide 1: 타이틀 ---
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "서을 슬세권 지수 분석 대시보드"
    subtitle.text = "공공데이터와 지리정보(GIS)를 활용한 내 집 앞 생활 인프라 분석\n발표자: 프로젝트 팀"

    # --- Slide 2: 도입 (배경) ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "왜 '슬세권'인가?"
    body = slide.shapes.placeholders[1].text_frame
    body.text = "현대인의 '근거리 생활권' 중시 트렌드"
    body.add_paragraph().text = "주관적인 '살기 좋은 동네'를 객관적 지표로 계량화"

    # --- Slide 3: 데이터 및 기술 스택 ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "데이터 및 기술 스택"
    body = slide.shapes.placeholders[1].text_frame
    body.text = "데이터: 서울시 공공데이터 (스타벅스, 지하철역, 병원 등 13종)"
    body.add_paragraph().text = "기술 스택: Python, Streamlit, Folium, Plotly"

    # --- Slide 4: 핵심 로직 (Engine) ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "핵심 로직 (Analysis Engine)"
    body = slide.shapes.placeholders[1].text_frame
    body.text = "거리 계산: geodesic 실측 거리 기반 필터링"
    body.add_paragraph().text = "지수 산출: 7개 카테고리 가중치 기반 100점 만점 설계"

    # --- Slide 5: 기능 시연 1 (지도/분석) ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "기능 시연 1: 실시간 지도 및 분석"
    body = slide.shapes.placeholders[1].text_frame
    body.text = "주소 검색 및 지도 클릭 연동"
    body.add_paragraph().text = "분석 반경 선택 (300m ~ 1500m)"

    # --- Slide 6: 기능 시연 2 (시각화) ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "기능 시연 2: 다각도 시각화"
    body = slide.shapes.placeholders[1].text_frame
    body.text = "레이더 차트: 인프라 균형도 분석"
    body.add_paragraph().text = "게이지 & 바 차트: 점수 및 평균 비교"

    # --- Slide 7: 개인화 기능 (Customizing) ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "개인화 및 리포트 기능"
    body = slide.shapes.placeholders[1].text_frame
    body.text = "사용자 맞춤형 가중치(Weight) 조정"
    body.add_paragraph().text = "분석 결과 HTML 보고서 저장 기능"

    # --- Slide 8: 결론 및 향후 과제 ---
    slide = prs.slides.add_slide(bullet_slide_layout)
    slide.shapes.title.text = "결론 및 향후 계획"
    body = slide.shapes.placeholders[1].text_frame
    body.text = "데이터 기반 주거 의사결정 지원"
    body.add_paragraph().text = "향후 실거래가 및 유동인구 데이터 연동"

    # 3. 파일 저장
    output_path = "docs/presentation.pptx"
    prs.save(output_path)
    print(f"PPT 생성 완료: {os.path.abspath(output_path)}")

if __name__ == "__main__":
    create_presentation()
